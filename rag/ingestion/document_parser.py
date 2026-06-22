"""
rag/ingestion/document_parser.py
Multi-format document parser: PDF, CSV, JSON, DOCX, XLSX, TXT, HTML, XML, PPTX
"""
import json
import csv
import io
import os
from pathlib import Path
from typing import List, Dict, Any
from dataclasses import dataclass, field
from loguru import logger


@dataclass
class ParsedChunk:
    content: str
    metadata: Dict[str, Any] = field(default_factory=dict)


class DocumentParser:
    """
    Parses various document formats into text chunks.
    Returns list of ParsedChunk objects.
    """

    PARSERS = {
        ".pdf":  "_parse_pdf",
        ".csv":  "_parse_csv",
        ".json": "_parse_json",
        ".txt":  "_parse_txt",
        ".docx": "_parse_docx",
        ".doc":  "_parse_docx",
        ".xlsx": "_parse_xlsx",
        ".xls":  "_parse_xlsx",
        ".pptx": "_parse_pptx",
        ".html": "_parse_html",
        ".htm":  "_parse_html",
        ".xml":  "_parse_xml",
    }

    def parse(self, file_path: str | Path) -> List[ParsedChunk]:
        """Dispatch to the correct parser based on file extension."""
        path = Path(file_path)
        ext  = path.suffix.lower()
        method_name = self.PARSERS.get(ext)

        if not method_name:
            raise ValueError(f"Unsupported file type: {ext}")

        logger.info(f"Parsing {path.name} as {ext}")
        parser = getattr(self, method_name)
        return parser(path)

    # ─── PDF ──────────────────────────────────────────────────────────────────
    def _parse_pdf(self, path: Path) -> List[ParsedChunk]:
        try:
            import pdfplumber
            chunks = []
            with pdfplumber.open(path) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    text = page.extract_text() or ""
                    if text.strip():
                        chunks.append(ParsedChunk(
                            content=text.strip(),
                            metadata={"page": page_num, "source": path.name, "type": "pdf"}
                        ))
                    # Also extract tables from PDF
                    tables = page.extract_tables()
                    for t_idx, table in enumerate(tables):
                        rows = [" | ".join(str(c) for c in row if c) for row in table if row]
                        table_text = "\n".join(rows)
                        if table_text.strip():
                            chunks.append(ParsedChunk(
                                content=table_text,
                                metadata={"page": page_num, "table": t_idx, "source": path.name, "type": "pdf_table"}
                            ))
            return chunks
        except Exception as e:
            logger.error(f"PDF parse error: {e}")
            raise

    # ─── CSV ──────────────────────────────────────────────────────────────────
    def _parse_csv(self, path: Path) -> List[ParsedChunk]:
        import pandas as pd
        chunks = []
        try:
            df = pd.read_csv(path, encoding="utf-8-sig")
        except UnicodeDecodeError:
            df = pd.read_csv(path, encoding="latin-1")

        # Header summary
        header_text = f"CSV File: {path.name}\nKolom: {', '.join(df.columns.tolist())}\nTotal baris: {len(df)}"
        chunks.append(ParsedChunk(
            content=header_text,
            metadata={"source": path.name, "type": "csv_header"}
        ))

        # Chunk rows (50 rows per chunk)
        chunk_size = 50
        for start in range(0, len(df), chunk_size):
            subset = df.iloc[start:start + chunk_size]
            text = subset.to_string(index=False)
            chunks.append(ParsedChunk(
                content=text,
                metadata={"source": path.name, "rows": f"{start}-{start+len(subset)}", "type": "csv"}
            ))
        return chunks

    # ─── JSON ─────────────────────────────────────────────────────────────────
    def _parse_json(self, path: Path) -> List[ParsedChunk]:
        with open(path, "r", encoding="utf-8") as f:
            data = json.load(f)
        return self._flatten_json(data, path.name)

    def _flatten_json(self, data: Any, source: str, prefix: str = "") -> List[ParsedChunk]:
        chunks = []
        if isinstance(data, list):
            for i, item in enumerate(data):
                sub = self._flatten_json(item, source, prefix=f"[{i}]")
                chunks.extend(sub)
        elif isinstance(data, dict):
            text = json.dumps(data, ensure_ascii=False, indent=2)
            if len(text) > 2000:
                # Split large dict into key groups
                items = list(data.items())
                for i in range(0, len(items), 10):
                    chunk_dict = dict(items[i:i+10])
                    chunks.append(ParsedChunk(
                        content=json.dumps(chunk_dict, ensure_ascii=False, indent=2),
                        metadata={"source": source, "type": "json", "keys": list(chunk_dict.keys())}
                    ))
            else:
                chunks.append(ParsedChunk(
                    content=text,
                    metadata={"source": source, "type": "json"}
                ))
        else:
            chunks.append(ParsedChunk(
                content=str(data),
                metadata={"source": source, "type": "json_value"}
            ))
        return chunks

    # ─── TXT ──────────────────────────────────────────────────────────────────
    def _parse_txt(self, path: Path) -> List[ParsedChunk]:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
        # Split on double newlines (paragraphs)
        paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
        return [
            ParsedChunk(content=p, metadata={"source": path.name, "type": "txt", "para": i})
            for i, p in enumerate(paragraphs)
        ]

    # ─── DOCX ─────────────────────────────────────────────────────────────────
    def _parse_docx(self, path: Path) -> List[ParsedChunk]:
        from docx import Document as DocxDocument
        doc = DocxDocument(path)
        chunks = []
        current_section = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                if current_section:
                    chunks.append(ParsedChunk(
                        content="\n".join(current_section),
                        metadata={"source": path.name, "type": "docx"}
                    ))
                    current_section = []
            else:
                current_section.append(text)
        if current_section:
            chunks.append(ParsedChunk(
                content="\n".join(current_section),
                metadata={"source": path.name, "type": "docx"}
            ))
        # Tables in docx
        for t_idx, table in enumerate(doc.tables):
            rows = [" | ".join(cell.text for cell in row.cells) for row in table.rows]
            chunks.append(ParsedChunk(
                content="\n".join(rows),
                metadata={"source": path.name, "type": "docx_table", "table": t_idx}
            ))
        return chunks

    # ─── XLSX ─────────────────────────────────────────────────────────────────
    def _parse_xlsx(self, path: Path) -> List[ParsedChunk]:
        import pandas as pd
        chunks = []
        xl = pd.ExcelFile(path)
        for sheet in xl.sheet_names:
            df = pd.read_excel(path, sheet_name=sheet)
            header = f"Sheet: {sheet} | Kolom: {', '.join(df.columns.astype(str).tolist())} | Baris: {len(df)}"
            chunks.append(ParsedChunk(
                content=header,
                metadata={"source": path.name, "sheet": sheet, "type": "xlsx_header"}
            ))
            for start in range(0, len(df), 50):
                subset = df.iloc[start:start+50]
                chunks.append(ParsedChunk(
                    content=subset.to_string(index=False),
                    metadata={"source": path.name, "sheet": sheet, "type": "xlsx",
                              "rows": f"{start}-{start+len(subset)}"}
                ))
        return chunks

    # ─── PPTX ─────────────────────────────────────────────────────────────────
    def _parse_pptx(self, path: Path) -> List[ParsedChunk]:
        from pptx import Presentation
        prs = Presentation(path)
        chunks = []
        for slide_num, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            if texts:
                chunks.append(ParsedChunk(
                    content="\n".join(texts),
                    metadata={"source": path.name, "slide": slide_num, "type": "pptx"}
                ))
        return chunks

    # ─── HTML ─────────────────────────────────────────────────────────────────
    def _parse_html(self, path: Path) -> List[ParsedChunk]:
        from bs4 import BeautifulSoup
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            soup = BeautifulSoup(f, "lxml")
        for tag in soup(["script", "style"]):
            tag.decompose()
        text = soup.get_text(separator="\n")
        paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
        return [
            ParsedChunk(content=p, metadata={"source": path.name, "type": "html"})
            for p in paragraphs
        ]

    # ─── XML ──────────────────────────────────────────────────────────────────
    def _parse_xml(self, path: Path) -> List[ParsedChunk]:
        from bs4 import BeautifulSoup
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            soup = BeautifulSoup(f, "xml")
        text = soup.get_text(separator="\n")
        paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
        return [
            ParsedChunk(content=p, metadata={"source": path.name, "type": "xml"})
            for p in paragraphs
        ]
